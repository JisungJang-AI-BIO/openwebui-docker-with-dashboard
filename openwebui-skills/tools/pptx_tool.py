"""
title: PPTX Presentation Tool
author: Internal Team
author_url: https://github.com/anthropics/skills
description: Read, create, edit, and convert PowerPoint presentations (.pptx).
    Use when user uploads a presentation or requests slide deck creation,
    editing, or conversion to PDF/images.
required_open_webui_version: 0.8.5
requirements: python-pptx, Pillow, pdf2image
version: 1.0.0
licence: MIT
"""

import asyncio
import base64
import io
import json
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Callable, Optional

from pydantic import BaseModel, Field


# =============================================================================
# EventEmitter Helper
# =============================================================================
class EventEmitter:
    def __init__(self, event_emitter: Callable[[dict], Any] = None, user: dict = None):
        self.event_emitter = event_emitter
        self.user = user or {}

    async def emit(self, description="Unknown State", status="in_progress", done=False):
        if self.event_emitter:
            await self.event_emitter(
                {"type": "status", "data": {"status": status, "description": description, "done": done}}
            )

    async def progress(self, description: str):
        await self.emit(description)

    async def success(self, description: str):
        await self.emit(description, "success", True)

    async def error(self, description: str):
        await self.emit(description, "error", True)

    async def send_file_link(self, file_path: str, filename: str, mime_type: str = None):
        """Upload file to OpenWebUI Files API and emit download link.
        Falls back to base64 data URI if the API is unavailable.
        """
        if not self.event_emitter or not os.path.exists(file_path):
            return
        if mime_type is None:
            ext = Path(filename).suffix.lower()
            mime_map = {
                ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ".pdf": "application/pdf", ".html": "text/html", ".txt": "text/plain",
                ".csv": "text/csv", ".odt": "application/vnd.oasis.opendocument.text",
                ".rtf": "application/rtf", ".png": "image/png",
                ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".gif": "image/gif",
            }
            mime_type = mime_map.get(ext, "application/octet-stream")

        is_image = mime_type.startswith("image/")
        url = await self._upload_to_openwebui(file_path, filename, mime_type)

        if url:
            if is_image:
                content = f"\n\n📎 **{filename}**\n\n![{filename}]({url})\n\n[Download {filename}]({url})\n"
            else:
                content = f"\n\n📎 **{filename}**\n\n[Download {filename}]({url})\n"
        else:
            with open(file_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("utf-8")
            data_uri = f"data:{mime_type};base64,{b64}"
            if is_image:
                content = f"\n\n📎 **{filename}**\n\n![{filename}]({data_uri})\n\n[Download {filename}]({data_uri})\n"
            else:
                content = f"\n\n📎 **{filename}**\n\n[Download {filename}]({data_uri})\n"

        await self.event_emitter({"type": "message", "data": {"content": content}})

    async def _upload_to_openwebui(self, file_path: str, filename: str, mime_type: str) -> Optional[str]:
        """Upload file via OpenWebUI's internal Files API. Returns download URL or None."""
        try:
            import httpx
            import jwt

            user_id = self.user.get("id", "")
            secret = os.environ.get("WEBUI_SECRET_KEY", "")
            if not user_id or not secret:
                return None

            token = jwt.encode({"id": user_id}, secret, algorithm="HS256")
            if isinstance(token, bytes):
                token = token.decode("utf-8")

            port = os.environ.get("PORT", "8080")
            with open(file_path, "rb") as f:
                resp = httpx.post(
                    f"http://localhost:{port}/api/v1/files/",
                    files={"file": (filename, f, mime_type)},
                    headers={"Authorization": f"Bearer {token}"},
                    timeout=30.0,
                )

            if resp.status_code in (200, 201):
                file_id = resp.json().get("id")
                if file_id:
                    return f"/api/v1/files/{file_id}/content"
            return None
        except Exception:
            return None


# =============================================================================
# Main Tool Class
# =============================================================================
class Tools:
    class Valves(BaseModel):
        TEMP_DIR: str = Field(default="/tmp/openwebui-pptx", description="Temporary file storage path")
        MAX_FILE_SIZE_MB: int = Field(default=50, description="Maximum file size in MB")
        SCRIPTS_DIR: str = Field(
            default="",
            description="Absolute path to Anthropic pptx scripts (e.g. /app/OpenWebUI-Skills/vendor/pptx). Leave empty for fallback mode.",
        )
        LIBREOFFICE_PATH: str = Field(default="soffice", description="LibreOffice binary path")
        DEFAULT_FONT: str = Field(default="Arial", description="Default font for new presentations")

    def __init__(self):
        self.valves = self.Valves()
        self.citation = False

    # =========================================================================
    # Tool Methods
    # =========================================================================

    async def read_pptx(
        self,
        mode: str = "text",
        __files__: list = None,
        __user__: dict = None,
        __event_emitter__: Callable[[dict], Any] = None,
    ) -> str:
        """
        Read and extract content from an uploaded PowerPoint file (.pptx).

        Modes:
        - "text": Plain text per slide
        - "structured": Detailed structure (layouts, shapes, positions, notes)

        :param mode: "text" or "structured"
        :return: Extracted slide content
        """
        emitter = EventEmitter(__event_emitter__, __user__)
        await emitter.progress("Reading PPTX...")

        if not __files__:
            await emitter.error("No file uploaded")
            return "Error: No file uploaded. Please upload a .pptx file."

        file_path = self._resolve_pptx_file(__files__)
        if not file_path:
            await emitter.error("No PPTX file found")
            return "Error: No .pptx file found."

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu

            prs = Presentation(file_path)
            parts = [f"## Presentation: {Path(file_path).name}\n"]
            parts.append(f"**Slides**: {len(prs.slides)}")
            w, h = prs.slide_width, prs.slide_height
            parts.append(f"**Size**: {w/914400:.1f}\" x {h/914400:.1f}\"\n")

            for i, slide in enumerate(prs.slides):
                slide_num = i + 1
                layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
                parts.append(f"### Slide {slide_num} (Layout: {layout_name})")

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                if mode == "structured":
                                    parts.append(f"  [{shape.shape_type}] {text}")
                                else:
                                    parts.append(text)
                    if shape.has_table:
                        table = shape.table
                        parts.append(f"\n**Table ({len(table.rows)}x{len(table.columns)}):**")
                        for r_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            parts.append("| " + " | ".join(cells) + " |")
                            if r_idx == 0:
                                parts.append("| " + " | ".join(["---"] * len(cells)) + " |")

                if mode == "structured" and slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"  **Notes**: {notes}")

                parts.append("")

            await emitter.success(f"Read {len(prs.slides)} slides")
            return "\n".join(parts)

        except Exception as e:
            await emitter.error(f"Read error: {str(e)}")
            return f"Error reading PPTX: {str(e)}"

    async def create_pptx(
        self,
        content: str,
        filename: str = "presentation.pptx",
        theme: str = "default",
        __user__: dict = None,
        __event_emitter__: Callable[[dict], Any] = None,
    ) -> str:
        """
        Create a new PowerPoint presentation.

        Content format — one slide per section separated by '---':
        ```
        # Slide Title
        Subtitle or body text

        - Bullet point 1
        - Bullet point 2

        | Col A | Col B |
        | Data 1 | Data 2 |
        ---
        ## Next Slide Title
        More content here
        ```

        - '#' = Title slide (large heading + optional subtitle)
        - '##' = Section heading slide
        - '###' = Content slide heading
        - '-' = Bullet points
        - '| |' = Table rows
        - '---' = Slide separator

        DESIGN RULES (avoid AI-generated look):
        - NEVER use accent lines under titles
        - Use asymmetric layouts when possible
        - Limit text per slide — if over 5 bullets, split
        - Use consistent color palette throughout

        :param content: Slide content specification
        :param filename: Output filename
        :param theme: Color theme (default, dark, minimal)
        :return: Status with download link
        """
        emitter = EventEmitter(__event_emitter__, __user__)
        await emitter.progress("Creating PPTX...")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Theme colors
            themes = {
                "default": {"bg": RGBColor(0xFF, 0xFF, 0xFF), "title": RGBColor(0x1A, 0x1A, 0x2E), "text": RGBColor(0x33, 0x33, 0x33), "accent": RGBColor(0x21, 0x96, 0xF3)},
                "dark": {"bg": RGBColor(0x1A, 0x1A, 0x2E), "title": RGBColor(0xFF, 0xFF, 0xFF), "text": RGBColor(0xCC, 0xCC, 0xCC), "accent": RGBColor(0x00, 0xBC, 0xD4)},
                "minimal": {"bg": RGBColor(0xFA, 0xFA, 0xFA), "title": RGBColor(0x21, 0x21, 0x21), "text": RGBColor(0x42, 0x42, 0x42), "accent": RGBColor(0xFF, 0x57, 0x22)},
            }
            colors = themes.get(theme, themes["default"])

            slides_content = content.split("---")

            for slide_text in slides_content:
                slide_text = slide_text.strip()
                if not slide_text:
                    continue
                self._build_slide(prs, slide_text, colors)

            temp_dir = self._ensure_temp_dir()
            output_path = os.path.join(temp_dir, filename)
            prs.save(output_path)

            await emitter.send_file_link(output_path, filename)
            await emitter.success(f"Created {filename}")

            file_size = os.path.getsize(output_path)
            return f"Created '{filename}' ({file_size:,} bytes, {len(prs.slides)} slides)."

        except Exception as e:
            await emitter.error(f"Creation error: {str(e)}")
            return f"Error creating PPTX: {str(e)}"

    async def edit_pptx(
        self,
        operation: str,
        parameters: str = "",
        __files__: list = None,
        __user__: dict = None,
        __event_emitter__: Callable[[dict], Any] = None,
    ) -> str:
        """
        Edit an uploaded PPTX file.

        Operations:
        - "replace_text": Replace text across all slides. parameters = "old_text|||new_text"
        - "delete_slide": Delete a slide by number. parameters = slide number (1-indexed)
        - "add_slide": Add a blank slide at end. parameters = "title|||body_text" (optional)

        :param operation: Operation name
        :param parameters: Operation-specific parameters
        :return: Edited file download link
        """
        emitter = EventEmitter(__event_emitter__, __user__)
        await emitter.progress(f"Editing PPTX: {operation}...")

        if not __files__:
            await emitter.error("No file uploaded")
            return "Error: No file uploaded."

        file_path = self._resolve_pptx_file(__files__)
        if not file_path:
            return "Error: No .pptx file found."

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            base_name = Path(file_path).stem

            if operation == "replace_text":
                parts = parameters.split("|||")
                if len(parts) != 2:
                    return "Error: format is 'old_text|||new_text'"
                old_text, new_text = parts
                count = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text)
                                        count += 1
                msg = f"Replaced {count} occurrence(s)."

            elif operation == "delete_slide":
                slide_num = int(parameters.strip())
                idx = slide_num - 1
                if 0 <= idx < len(prs.slides):
                    rId = prs.slides._sldIdLst[idx].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[idx]
                    msg = f"Deleted slide {slide_num}."
                else:
                    return f"Error: Slide {slide_num} out of range (1-{len(prs.slides)})."

            elif operation == "add_slide":
                layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                if parameters:
                    parts = parameters.split("|||")
                    if slide.placeholders:
                        for i, ph in enumerate(slide.placeholders):
                            if i < len(parts):
                                ph.text = parts[i]
                msg = f"Added slide {len(prs.slides)}."

            else:
                return f"Error: Unknown operation '{operation}'. Available: replace_text, delete_slide, add_slide"

            temp_dir = self._ensure_temp_dir()
            output_filename = f"{base_name}_edited.pptx"
            output_path = os.path.join(temp_dir, output_filename)
            prs.save(output_path)

            await emitter.send_file_link(output_path, output_filename)
            await emitter.success(msg)
            return msg

        except Exception as e:
            await emitter.error(f"Edit error: {str(e)}")
            return f"Error: {str(e)}"

    async def convert_pptx(
        self,
        target_format: str = "pdf",
        __files__: list = None,
        __user__: dict = None,
        __event_emitter__: Callable[[dict], Any] = None,
    ) -> str:
        """
        Convert an uploaded PPTX to another format using LibreOffice.
        Supported: pdf, png, jpg.

        :param target_format: "pdf", "png", or "jpg"
        :return: Converted file download link
        """
        emitter = EventEmitter(__event_emitter__, __user__)
        await emitter.progress(f"Converting to {target_format}...")

        if not __files__:
            await emitter.error("No file uploaded")
            return "Error: No file uploaded."

        file_path = self._resolve_pptx_file(__files__)
        if not file_path:
            return "Error: No .pptx file found."

        try:
            temp_dir = self._ensure_temp_dir()
            fmt = target_format.lower().strip(".")

            if fmt in ("png", "jpg", "jpeg"):
                return await self._convert_to_images(file_path, temp_dir, fmt, emitter)
            else:
                return await self._convert_libreoffice(file_path, temp_dir, fmt, emitter)

        except Exception as e:
            await emitter.error(f"Conversion error: {str(e)}")
            return f"Error: {str(e)}"

    async def export_slides_as_images(
        self,
        __files__: list = None,
        __user__: dict = None,
        __event_emitter__: Callable[[dict], Any] = None,
    ) -> str:
        """
        Export all slides as individual PNG images.
        Pipeline: pptx → pdf (LibreOffice) → png (pdftoppm).

        :return: Image download links for each slide
        """
        emitter = EventEmitter(__event_emitter__, __user__)

        if not __files__:
            await emitter.error("No file uploaded")
            return "Error: No file uploaded."

        file_path = self._resolve_pptx_file(__files__)
        if not file_path:
            return "Error: No .pptx file found."

        return await self._convert_to_images(file_path, self._ensure_temp_dir(), "png", emitter)

    # =========================================================================
    # Internal helpers
    # =========================================================================

    def _resolve_pptx_file(self, files: list) -> Optional[str]:
        if not files:
            return None
        for f in files:
            info = f if isinstance(f, dict) else {"path": f}
            name = info.get("filename", info.get("name", str(info.get("path", ""))))
            if not name.lower().endswith((".pptx", ".ppt")):
                continue
            for key in ("path", "file_path", "url", "id"):
                path = info.get(key)
                if path and isinstance(path, str):
                    for c in [path, f"/app/backend/data/uploads/{path}", f"/app/backend/data/cache/{path}"]:
                        if os.path.exists(c):
                            return c
                    return path
        if files:
            f = files[0]
            if isinstance(f, dict):
                for key in ("path", "file_path"):
                    if key in f:
                        return f[key]
            elif isinstance(f, str):
                return f
        return None

    def _ensure_temp_dir(self) -> str:
        os.makedirs(self.valves.TEMP_DIR, exist_ok=True)
        return self.valves.TEMP_DIR

    def _build_slide(self, prs, slide_text: str, colors: dict):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        lines = slide_text.strip().split("\n")
        if not lines:
            return

        # Determine slide type from first line
        first = lines[0].strip()
        is_title = first.startswith("# ") and not first.startswith("## ")

        layout_idx = 0 if is_title else 1
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 0
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = colors["bg"]

        title_text = ""
        body_lines = []
        table_buf = []

        for line in lines:
            s = line.strip()
            if s.startswith("# ") and not title_text:
                title_text = s.lstrip("#").strip()
            elif s.startswith("|") and s.endswith("|"):
                inner = s[1:-1].strip()
                if not all(c in "-| :" for c in inner):
                    table_buf.append(s)
            else:
                if table_buf:
                    body_lines.append(("table", table_buf))
                    table_buf = []
                if s:
                    body_lines.append(("text", s))
        if table_buf:
            body_lines.append(("table", table_buf))

        # Set title
        if slide.placeholders and title_text:
            ph = slide.placeholders[0]
            ph.text = title_text
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.color.rgb = colors["title"]
                run.font.size = Pt(36 if is_title else 28)
                run.font.bold = True

        # Set body content
        if len(slide.placeholders) > 1:
            body_ph = slide.placeholders[1]
            tf = body_ph.text_frame
            tf.clear()
            first_para = True
            for item_type, item in body_lines:
                if item_type == "text":
                    s = item
                    if s.startswith("# "):
                        continue  # already used as title
                    p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                    first_para = False
                    if s.startswith("- ") or s.startswith("* "):
                        p.text = s[2:]
                        p.level = 0
                    elif s.startswith("  - "):
                        p.text = s[4:]
                        p.level = 1
                    else:
                        p.text = s
                    for run in p.runs:
                        run.font.color.rgb = colors["text"]
                        run.font.size = Pt(18)

    async def _convert_libreoffice(self, file_path, temp_dir, fmt, emitter) -> str:
        soffice = self.valves.LIBREOFFICE_PATH
        input_copy = os.path.join(temp_dir, Path(file_path).name)
        shutil.copy2(file_path, input_copy)
        env = os.environ.copy()
        env["SAL_USE_VCLPLUGIN"] = "svp"
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", fmt, "--outdir", temp_dir, input_copy],
                capture_output=True, text=True, timeout=120, env=env,
            )
        except FileNotFoundError:
            return f"Error: LibreOffice not found at '{soffice}'."
        except subprocess.TimeoutExpired:
            return "Error: Conversion timed out."

        output_path = os.path.join(temp_dir, f"{Path(file_path).stem}.{fmt}")
        if os.path.exists(output_path):
            await emitter.send_file_link(output_path, Path(output_path).name)
            await emitter.success(f"Converted to {fmt}")
            return f"Converted to {fmt} ({os.path.getsize(output_path):,} bytes)."
        return f"Error: Conversion failed. {result.stderr}"

    async def _convert_to_images(self, file_path, temp_dir, fmt, emitter) -> str:
        # Step 1: pptx → pdf
        soffice = self.valves.LIBREOFFICE_PATH
        input_copy = os.path.join(temp_dir, Path(file_path).name)
        shutil.copy2(file_path, input_copy)
        env = os.environ.copy()
        env["SAL_USE_VCLPLUGIN"] = "svp"

        await emitter.progress("Converting to PDF first...")
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, input_copy],
                capture_output=True, text=True, timeout=120, env=env,
            )
        except FileNotFoundError:
            return f"Error: LibreOffice not found at '{soffice}'."

        pdf_path = os.path.join(temp_dir, f"{Path(file_path).stem}.pdf")
        if not os.path.exists(pdf_path):
            return "Error: PDF conversion failed."

        # Step 2: pdf → images
        await emitter.progress("Converting PDF to images...")
        img_fmt = fmt.replace("jpg", "jpeg")
        try:
            from pdf2image import convert_from_path

            images = convert_from_path(pdf_path, dpi=150, fmt=img_fmt)
            saved = []
            for i, img in enumerate(images):
                fname = f"{Path(file_path).stem}_slide{i+1}.{fmt}"
                fpath = os.path.join(temp_dir, fname)
                img.save(fpath)
                saved.append((fpath, fname))

            for fpath, fname in saved:
                await emitter.send_file_link(fpath, fname)

            await emitter.success(f"Exported {len(saved)} slides as {fmt}")
            return f"Exported {len(saved)} slides as {fmt} images."

        except Exception as e:
            return f"Error converting to images: {str(e)}"
